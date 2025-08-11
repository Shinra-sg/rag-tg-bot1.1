#!/usr/bin/env python3
"""
Парсер для PowerPoint файлов (.pptx, .ppt)
"""

import sys
import json
import os
from pptx import Presentation
from datetime import datetime

def parse_powerpoint(file_path):
    """Парсит PowerPoint файл и возвращает текст и метаданные"""
    try:
        # Загружаем презентацию
        prs = Presentation(file_path)
        
        content_parts = []
        
        # Обрабатываем каждый слайд
        for slide_num, slide in enumerate(prs.slides, 1):
            content_parts.append(f"=== Слайд {slide_num} ===")
            
            # Извлекаем текст из фигур на слайде
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content_parts.append(shape.text)
            
            # Извлекаем текст из таблиц
            for shape in slide.shapes:
                if shape.has_table:
                    content_parts.append("--- Таблица ---")
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        if row_text.strip():
                            content_parts.append(row_text)
                    content_parts.append("")
            
            content_parts.append("")  # Пустая строка между слайдами
        
        content = "\n".join(content_parts)
        
        # Получаем метаданные
        metadata = {
            "title": os.path.basename(file_path),
            "author": None,
            "slides": len(prs.slides),
            "created": None,
            "modified": None
        }
        
        # Пытаемся получить метаданные из свойств презентации
        if hasattr(prs, 'core_properties'):
            if prs.core_properties.author:
                metadata["author"] = prs.core_properties.author
            if prs.core_properties.created:
                metadata["created"] = prs.core_properties.created.isoformat()
            if prs.core_properties.modified:
                metadata["modified"] = prs.core_properties.modified.isoformat()
        
        return {
            "content": content,
            "metadata": metadata
        }
        
    except Exception as e:
        return {
            "content": f"Ошибка парсинга PowerPoint файла: {str(e)}",
            "metadata": {"title": os.path.basename(file_path)}
        }

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(json.dumps({"error": "Необходимо указать путь к файлу"}))
        sys.exit(1)
    
    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(json.dumps({"error": f"Файл не найден: {file_path}"}))
        sys.exit(1)
    
    result = parse_powerpoint(file_path)
    print(json.dumps(result, ensure_ascii=False, default=str)) 